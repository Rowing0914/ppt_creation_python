from pptx.util import Inches
from pptx import Presentation

import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import sqlite3
import csv

df = pd.read_csv('../data/sample.csv')
print(df.iloc[:, 1])

plt.scatter(df.iloc[:, 1], df.iloc[:, 2])

filename = "../image/output.png"
plt.savefig(filename)

def open_PowerPoint_Presentation(oldFileName, newFileName, img, left, top):
    prs = Presentation(oldFileName)
    slide = prs.slides[1]
    pic = slide.shapes.add_picture(img, left, top)
    prs.save(newFileName)

open_PowerPoint_Presentation('../ppt/sample.pptx', '../ppt/NewTemplate.pptx',
                             '../image/output.png', Inches(1), Inches(1))